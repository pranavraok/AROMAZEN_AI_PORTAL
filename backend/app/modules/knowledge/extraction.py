from itertools import chain, islice
from pathlib import Path

from docx import Document as WordDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


class ExtractionError(Exception):
    """Raised when a supported document cannot be read locally."""


def _clean(parts: list[str]) -> str:
    return "\n".join(part.strip() for part in parts if part and part.strip())


def _cell_text(value) -> str:
    return "" if value is None else " ".join(str(value).split())


def _word_paragraph_text(paragraph) -> str:
    """Preserve degree symbols represented as a superscript zero in Word files."""
    return "".join(
        "°" if run.font.superscript and run.text == "0" else run.text
        for run in paragraph.runs
    )


def _word_cell_text(cell) -> str:
    return "\n".join(_word_paragraph_text(paragraph) for paragraph in cell.paragraphs)


def _spreadsheet_text(workbook) -> str:
    """Represent sheets as labeled records so semantic search understands each value."""
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"Worksheet: {sheet.title}")
        rows = sheet.iter_rows(values_only=True)
        buffered_rows = list(islice(rows, 25))
        header_index = None
        if buffered_rows:
            scores = [
                sum(isinstance(value, str) and bool(value.strip()) for value in row)
                for row in buffered_rows
            ]
            best_index = max(range(len(scores)), key=scores.__getitem__)
            if scores[best_index] >= 2:
                header_index = best_index
        headers = (
            [_cell_text(value) for value in buffered_rows[header_index]]
            if header_index is not None
            else []
        )
        for row_index, row in enumerate(chain(buffered_rows, rows)):
            values = [_cell_text(value) for value in row]
            if not any(values):
                continue
            if header_index is None or row_index <= header_index:
                parts.append(" | ".join(value for value in values if value))
                continue
            labeled_values = [
                f"{headers[index] or f'Column {index + 1}'}: {value}"
                for index, value in enumerate(values)
                if value
            ]
            parts.append(f"Record {row_index + 1}: " + " | ".join(labeled_values))
    return _clean(parts)


def extract_text(file_path: Path, extension: str) -> str:
    """Extract readable text locally. No external AI service is used."""
    try:
        if extension == ".pdf":
            return _clean([page.extract_text() or "" for page in PdfReader(str(file_path)).pages])
        if extension == ".docx":
            document = WordDocument(str(file_path))
            return _clean(
                [_word_paragraph_text(paragraph) for paragraph in document.paragraphs]
                + [" | ".join(_word_cell_text(cell) for cell in row.cells) for table in document.tables for row in table.rows]
            )
        if extension == ".xlsx":
            workbook = load_workbook(str(file_path), read_only=True, data_only=True)
            return _spreadsheet_text(workbook)
        if extension == ".pptx":
            presentation = Presentation(str(file_path))
            parts = []
            for number, slide in enumerate(presentation.slides, start=1):
                parts.append(f"Slide {number}")
                parts.extend(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            return _clean(parts)
    except Exception as exc:
        raise ExtractionError("This file could not be read. Please confirm it is a valid, unprotected document.") from exc
    raise ExtractionError("This file type is not supported.")
